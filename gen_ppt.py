from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# White background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Title ──
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "量化产品线一览"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.alignment = PP_ALIGN.LEFT

def add_section(text, left, top, width, height):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x52, 0x8B)
    p.alignment = PP_ALIGN.LEFT
    return tf

def add_body(lines, left, top, width, height, font_size=18):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT
    return tf

# ── Column Layout (4 columns) ──
col_w = 2.9
col_gap = 0.35
col_start = 0.55

# Col 0: 沪深300指增
add_section("沪深300指增", col_start + 0 * (col_w + col_gap), 1.2, col_w, 0.5)
add_body([
    "■ 潮汐300指增",
    "  仓位：60%-140%",
    "  动态换手率：80-200倍",
    "  持股：700-1000只",
    "  中证300成分股≥50%",
    "  市值因子暴露<0.1",
], col_start + 0 * (col_w + col_gap), 1.75, col_w, 2.5, 18)

# Col 1: 中证1000指增
add_section("中证1000指增", col_start + 1 * (col_w + col_gap), 1.2, col_w, 0.5)
add_body([
    "■ 潮汐1000指增",
    "  仓位：60%-140%",
    "  动态换手率：80-200倍",
    "  持股：700-1000只",
    "  中证1000成分股≥50%",
    "  市值/非线性市值",
    "  因子暴露<0.3",
], col_start + 1 * (col_w + col_gap), 1.75, col_w, 2.8, 18)

# Col 2: 择时选股系列（上部分）
add_section("择时选股系列", col_start + 2.15 * (col_w + col_gap), 1.2, col_w, 0.5)
add_body([
    "■ 择时安渡",
    "  仓位：0%-100%",
    "  换手率：80-200倍",
    "  持股：500-1000只",
    "  全市场选股，剔除ST、",
    "  流动性差及市值后500",
], col_start + 2.15 * (col_w + col_gap), 1.75, col_w, 2.0, 18)

# Col 2 continued: 择时逐浪
add_body([
    "■ 择时逐浪",
    "  仓位：0-140%",
    "  换手率：80-200倍",
    "  持股：500-1000只",
    "  全市场选股，剔除ST、",
    "  流动性差及市值后500",
], col_start + 2.15 * (col_w + col_gap), 3.9, col_w, 2.0, 18)

# Col 2 continued: 双创择时
add_body([
    "■ 双创择时",
    "  仓位：0-140%",
    "  换手率：80-200倍",
    "  持股：200-300只",
    "  科创板和创业板成分股",
], col_start + 2.15 * (col_w + col_gap), 5.7, col_w, 1.8, 18)

# Col 3: 量化中性系列
add_section("量化中性系列", col_start + 3.3 * (col_w + col_gap), 1.2, col_w, 0.5)
add_body([
    "■ 择时中性",
    "  股票敞口：-10%~10%",
    "  通过择时信号决定",
    "  多头端300/500/1000权重",
    "  每15分钟比较远近合约",
    "  并进行移仓调整",
], col_start + 3.3 * (col_w + col_gap), 1.75, col_w, 2.5, 18)

# ── Separator lines between columns ──
for i in range(3):
    x = col_start + (i + 1) * (col_w + col_gap) - col_gap / 2 - 0.05
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(x), Inches(1.5), Inches(0.01), Inches(5.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    line.line.fill.background()

output_path = "/Users/wangqiuting/macro-dashboard/量化产品线一览.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
